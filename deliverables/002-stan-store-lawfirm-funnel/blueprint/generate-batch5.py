#!/usr/bin/env python3
"""
BATCH 5: Slides 21-25
- Slide 21: Section Break (Faceless Marketing Engine)
- Slide 22: Floating-Cards (Marketing automation components)
- Slide 23: Multi-Card (Content pipeline architecture)
- Slide 24: Section Break (Integration Architecture)
- Slide 25: Bold-Diagonal (System architecture overview)
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

BRAND_BG = "0a0a14"
BRAND_BG_ALT = "12121f"
BRAND_TEXT = "f5f5f5"
BRAND_TEXT_SECONDARY = "b8b8c8"
BRAND_ACCENT = "00ffff"
BRAND_ACCENT_SECONDARY = "ff00ff"
BRAND_CARD_BG = "1a1a2e"
BRAND_HEADING_FONT = "Playfair Display"
BRAND_BODY_FONT = "Inter"

def create_batch5():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 21: Section Break - Faceless Marketing Engine =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    num_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(4.0), Inches(2.0))
    p = num_box.text_frame.paragraphs[0]
    p.text = "07"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.8), Inches(3.0), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(10.0), Inches(1.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Faceless Marketing Engine"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(10.0), Inches(0.8))
    p = sub_box.text_frame.paragraphs[0]
    p.text = "Content that generates leads while you practice law"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # ===== SLIDE 22: Floating-Cards - Marketing Components =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = "The Marketing System Components"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    cards = [
        ("SEO Content Engine", "Practice-area specific articles targeting high-intent search terms in your geography. Automated research, outline, and draft pipeline."),
        ("Social Media Automation", "Faceless video content from your blog posts. LinkedIn thought leadership. YouTube shorts. Zero on-camera time required."),
        ("Review & Reputation", "Automated review requests post-case-close. Google Business Profile optimization. Reputation monitoring and response."),
    ]
    card_positions = [(0.5, 1.8), (4.6, 2.2), (8.7, 1.8)]
    colors = [BRAND_ACCENT, BRAND_ACCENT_SECONDARY, "00d4d4"]
    for i, ((title, desc), (x, y)) in enumerate(zip(cards, card_positions)):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.0), Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
        card.line.color.rgb = hex_to_rgb(colors[i])
        card.line.width = Pt(1.5)

        t_box = slide.shapes.add_textbox(Inches(x + 0.4), Inches(y + 0.5), Inches(3.2), Inches(0.6))
        p = t_box.text_frame.paragraphs[0]
        p.text = title
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(colors[i])

        d_box = slide.shapes.add_textbox(Inches(x + 0.4), Inches(y + 1.4), Inches(3.2), Inches(2.5))
        d_box.text_frame.word_wrap = True
        p = d_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(15)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # ===== SLIDE 23: Multi-Card - Content Pipeline =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = "The Content Pipeline"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    stages = [
        ("Research", "Keyword + competitor\nanalysis for your\npractice area + city"),
        ("Draft", "AI-assisted long-form\ncontent with attorney\nreview and approval"),
        ("Publish", "Blog, YouTube, LinkedIn\nautomated from single\nsource article"),
        ("Convert", "Every piece links to\nintake form — content\nis the top of funnel"),
    ]
    for i, (title, desc) in enumerate(stages):
        x = 0.5 + (i * 3.2)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(2.8), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
        card.line.color.rgb = hex_to_rgb(BRAND_ACCENT)
        card.line.width = Pt(1)

        n_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(2.3), Inches(1.0), Inches(0.6))
        p = n_box.text_frame.paragraphs[0]
        p.text = str(i + 1)
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

        t_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(3.0), Inches(2.2), Inches(0.5))
        p = t_box.text_frame.paragraphs[0]
        p.text = title
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

        d_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(3.6), Inches(2.2), Inches(1.5))
        d_box.text_frame.word_wrap = True
        p = d_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

        if i < 3:
            arrow = slide.shapes.add_textbox(Inches(x + 2.9), Inches(3.3), Inches(0.4), Inches(0.5))
            p = arrow.text_frame.paragraphs[0]
            p.text = "\u2192"
            p.font.size = Pt(28)
            p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    # ===== SLIDE 24: Section Break - Integration Architecture =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    num_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(4.0), Inches(2.0))
    p = num_box.text_frame.paragraphs[0]
    p.text = "08"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.8), Inches(3.0), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(10.0), Inches(1.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Integration Architecture"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(10.0), Inches(0.8))
    p = sub_box.text_frame.paragraphs[0]
    p.text = "How every system connects into one operating layer"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # ===== SLIDE 25: Bold-Diagonal - System Architecture =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Diagonal accent shape
    diag = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(0), Inches(0), Inches(6.0), Inches(7.5))
    diag.fill.solid()
    diag.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT)
    diag.line.fill.background()

    # Title on left
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(4.5), Inches(1.5))
    title_box.text_frame.word_wrap = True
    p = title_box.text_frame.paragraphs[0]
    p.text = "The Full System Architecture"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Architecture layers on right
    layers = [
        ("Client Layer", "Web forms, portals, e-sign, payment"),
        ("Automation Layer", "n8n workflows, Zapier, webhooks"),
        ("Data Layer", "Airtable, CRM, case management"),
        ("Intelligence Layer", "AI drafting, analytics, reporting"),
        ("Communication Layer", "Email, SMS, client portal, Slack"),
    ]
    for i, (layer, desc) in enumerate(layers):
        y = 0.8 + (i * 1.3)
        # Layer card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(y), Inches(6.0), Inches(1.1))
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
        card.line.color.rgb = hex_to_rgb(BRAND_ACCENT if i % 2 == 0 else BRAND_ACCENT_SECONDARY)
        card.line.width = Pt(1)

        l_box = slide.shapes.add_textbox(Inches(6.8), Inches(y + 0.1), Inches(2.5), Inches(0.4))
        p = l_box.text_frame.paragraphs[0]
        p.text = layer
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_ACCENT if i % 2 == 0 else BRAND_ACCENT_SECONDARY)

        d_box = slide.shapes.add_textbox(Inches(6.8), Inches(y + 0.5), Inches(5.5), Inches(0.4))
        p = d_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(13)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    return prs

if __name__ == "__main__":
    prs = create_batch5()
    prs.save("part5.pptx")
    print("✓ Batch 5 complete: Slides 21-25 saved to part5.pptx")

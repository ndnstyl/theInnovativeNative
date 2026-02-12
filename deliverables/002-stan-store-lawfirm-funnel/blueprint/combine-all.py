#!/usr/bin/env python3
"""Combine all 6 batch PPTX files into final presentation."""

import copy
import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

BRAND_BG = "0a0a14"

def main():
    # Start with part1 as base
    combined = Presentation("part1.pptx")
    combined.slide_width = Inches(13.333)
    combined.slide_height = Inches(7.5)
    blank_layout = combined.slide_layouts[6]

    # Append parts 2-6
    for part_num in range(2, 7):
        filename = f"part{part_num}.pptx"
        if not os.path.exists(filename):
            print(f"Warning: {filename} not found, skipping")
            continue

        part = Presentation(filename)
        for slide in part.slides:
            new_slide = combined.slides.add_slide(blank_layout)
            # CRITICAL: Set background explicitly on every combined slide
            new_slide.background.fill.solid()
            new_slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)
            # Copy all shapes
            for shape in slide.shapes:
                el = copy.deepcopy(shape.element)
                new_slide.shapes._spTree.append(el)

        print(f"✓ Added {filename} ({len(part.slides)} slides)")

    # Save final
    output = "automated-lawfirm-blueprint.pptx"
    combined.save(output)
    print(f"\n✓ Final presentation saved: {output}")
    print(f"  Total slides: {len(combined.slides)}")
    print(f"  File size: {os.path.getsize(output) / 1024:.1f} KB")

    # Clean up part files
    for i in range(1, 7):
        f = f"part{i}.pptx"
        if os.path.exists(f):
            os.remove(f)
            print(f"  Deleted {f}")

    print("\n✓ All part files cleaned up. Final PPTX ready.")

if __name__ == "__main__":
    main()

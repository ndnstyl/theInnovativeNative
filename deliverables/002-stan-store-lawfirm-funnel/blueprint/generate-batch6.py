#!/usr/bin/env python3
"""
BATCH 6: Slides 26-30
- Slide 26: Corner-Anchor (Integration stack layers)
- Slide 27: Section Break (Implementation Roadmap)
- Slide 28: Giant-Focus (90 Days to Automation)
- Slide 29: Stats (ROI projections)
- Slide 30: Closing Slide (CTA)
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

BRAND_BG = "0a0a14"
BRAND_BG_ALT = "12121f"
BRAND_TEXT = "f5f5f5"
BRAND_TEXT_SECONDARY = "b8b8c8"
BRAND_ACCENT = "00ffff"
BRAND_ACCENT_SECONDARY = "ff00ff"
BRAND_CARD_BG = "1a1a2e"
BRAND_HEADING_FONT = "Playfair Display"
BRAND_BODY_FONT = "Inter"

def create_batch6():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 26: Corner-Anchor - Integration Stack =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Corner accent block
    corner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(3.0))
    corner.fill.solid()
    corner.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    corner.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.5), Inches(11.0), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = "The Recommended Tech Stack"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    stacks = [
        ("Automation Hub", "n8n (self-hosted)", "Central workflow orchestration — connects every tool"),
        ("Data & CRM", "Airtable + Clio/MyCase", "Structured data layer + legal practice management"),
        ("Communication", "Twilio + SendGrid", "SMS and email automation for client touchpoints"),
        ("Documents", "Pandadoc + Google Docs", "Template-driven assembly with e-signature"),
        ("Marketing", "WordPress + Remotion", "SEO content + faceless video production"),
        ("Scheduling", "Cal.com", "Client-facing booking with automated reminders"),
    ]
    for i, (name, tool, desc) in enumerate(stacks):
        col = i % 2
        row = i // 2
        x = 1.0 + (col * 6.0)
        y = 1.8 + (row * 1.8)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.5), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
        card.line.color.rgb = hex_to_rgb(BRAND_ACCENT if col == 0 else BRAND_ACCENT_SECONDARY)
        card.line.width = Pt(1)

        n_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.2), Inches(2.5), Inches(0.4))
        p = n_box.text_frame.paragraphs[0]
        p.text = name
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_ACCENT if col == 0 else BRAND_ACCENT_SECONDARY)

        t_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.6), Inches(4.8), Inches(0.3))
        p = t_box.text_frame.paragraphs[0]
        p.text = tool
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

        d_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.95), Inches(4.8), Inches(0.4))
        d_box.text_frame.word_wrap = True
        p = d_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(12)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # ===== SLIDE 27: Section Break - Implementation Roadmap =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    num_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(4.0), Inches(2.0))
    p = num_box.text_frame.paragraphs[0]
    p.text = "09"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.8), Inches(3.0), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(10.0), Inches(1.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Implementation Roadmap"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(10.0), Inches(0.8))
    p = sub_box.text_frame.paragraphs[0]
    p.text = "A phased approach that starts generating ROI in 30 days"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # ===== SLIDE 28: Giant-Focus - 90 Days to Automation =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Giant text
    giant_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.0), Inches(2.0))
    giant_box.text_frame.word_wrap = True
    p = giant_box.text_frame.paragraphs[0]
    p.text = "90 Days"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(96)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.0), Inches(0.7))
    p = sub_box.text_frame.paragraphs[0]
    p.text = "From manual chaos to automated operations"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(28)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.CENTER

    phases = [
        ("Days 1-30", "Foundation", "Intake automation + time tracking.\nImmediate ROI from lead capture\nand billing recovery.", BRAND_ACCENT),
        ("Days 31-60", "Expansion", "Document assembly + client comms.\nReduce associate hours on\nrepetitive work by 60%.", BRAND_ACCENT_SECONDARY),
        ("Days 61-90", "Scale", "Marketing engine + full integration.\nSelf-sustaining lead generation\nand operational automation.", "00d4d4"),
    ]
    for i, (timeline, phase, desc, color) in enumerate(phases):
        x = 0.5 + (i * 4.2)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.8), Inches(3.8), Inches(3.2))
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
        card.line.color.rgb = hex_to_rgb(color)
        card.line.width = Pt(1.5)

        tl_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(4.0), Inches(3.2), Inches(0.4))
        p = tl_box.text_frame.paragraphs[0]
        p.text = timeline
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = hex_to_rgb(color)

        ph_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(4.4), Inches(3.2), Inches(0.5))
        p = ph_box.text_frame.paragraphs[0]
        p.text = phase
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

        d_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(5.0), Inches(3.2), Inches(1.8))
        d_box.text_frame.word_wrap = True
        p = d_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # ===== SLIDE 29: Stats - ROI Projections =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Projected Annual ROI by Practice Area"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    stats = [
        ("$120K+", "Bankruptcy Solo Practice", "Recovered billable time + faster petition filing + automated creditor communications"),
        ("$150K+", "Criminal Defense Firm", "Lead capture speed + document assembly + automated case status updates"),
        ("$100K+", "Administrative Law Practice", "Deadline compliance + document automation + agency response efficiency"),
    ]
    for i, (value, label, desc) in enumerate(stats):
        y = 2.0 + (i * 1.7)
        v_box = slide.shapes.add_textbox(Inches(1.0), Inches(y), Inches(3.0), Inches(0.7))
        p = v_box.text_frame.paragraphs[0]
        p.text = value
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(48)
        p.font.bold = True
        colors = [BRAND_ACCENT, BRAND_ACCENT_SECONDARY, "00d4d4"]
        p.font.color.rgb = hex_to_rgb(colors[i])

        l_box = slide.shapes.add_textbox(Inches(4.5), Inches(y), Inches(7.5), Inches(0.5))
        p = l_box.text_frame.paragraphs[0]
        p.text = label
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

        d_box = slide.shapes.add_textbox(Inches(4.5), Inches(y + 0.5), Inches(7.5), Inches(0.5))
        d_box.text_frame.word_wrap = True
        p = d_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(13)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # ===== SLIDE 30: Closing Slide - CTA =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Bottom accent block
    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.0), Inches(13.333), Inches(2.5))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT)
    bottom.line.fill.background()

    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    bar.line.fill.background()

    # Headline
    h_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.0), Inches(1.5))
    h_box.text_frame.word_wrap = True
    p = h_box.text_frame.paragraphs[0]
    p.text = "Ready to Stop Losing $100K/Year?"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.CENTER

    # Subtext
    sub_box = slide.shapes.add_textbox(Inches(2.0), Inches(3.2), Inches(9.0), Inches(0.8))
    sub_box.text_frame.word_wrap = True
    p = sub_box.text_frame.paragraphs[0]
    p.text = "Book a free 30-minute diagnostic call. We'll review your scorecard results and map the first 90 days."
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
    p.alignment = PP_ALIGN.CENTER

    # CTAs
    ctas = [
        "Book Your Diagnostic Call → cal.com/theinnovativenative",
        "Download the ROI Calculator → included with this blueprint",
        "Take the Free Scorecard → stan.store/theinnovativenative",
    ]
    for i, cta in enumerate(ctas):
        y = 5.3 + (i * 0.6)
        c_box = slide.shapes.add_textbox(Inches(2.0), Inches(y), Inches(9.0), Inches(0.5))
        p = c_box.text_frame.paragraphs[0]
        p.text = cta
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(16)
        p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
        p.alignment = PP_ALIGN.CENTER

    # Brand
    brand_box = slide.shapes.add_textbox(Inches(3.0), Inches(4.2), Inches(7.0), Inches(0.5))
    p = brand_box.text_frame.paragraphs[0]
    p.text = "The Innovative Native LLC — theinnovativenative.com"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(14)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
    p.alignment = PP_ALIGN.CENTER

    return prs

if __name__ == "__main__":
    prs = create_batch6()
    prs.save("part6.pptx")
    print("✓ Batch 6 complete: Slides 26-30 saved to part6.pptx")

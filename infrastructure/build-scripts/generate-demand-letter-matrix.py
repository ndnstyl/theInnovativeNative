"""
generate-demand-letter-matrix.py
Generates demand-letter-matrix.pptx (6 slides) matching the existing PI lead-magnet PPTX style.
Source: projects/002-stan-store-lawfirm-funnel/content/lead-magnets/demand-letter-matrix.md
Output: projects/website/public/assets/lead-magnets/demand-letter-matrix.pptx

Style reference: 442-intake-math-pi.pptx
  - Background: #0A0A14
  - Accent bar color: #D4A853 (gold)
  - Stat highlight color: #00BCD4 (cyan)
  - Title font: Playfair Display, 42pt, bold, gold
  - Body font: Inter
  - Secondary text color: #B8B8C8
  - Table header fill: #D4A853, text #0A0A14
  - Table data rows: #121224 fill, text #F5F5F5
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import lxml.etree as etree
import copy

# ── Color palette ────────────────────────────────────────────────────────────
BG          = RGBColor(0x0A, 0x0A, 0x14)
ALT_BG      = RGBColor(0x12, 0x12, 0x24)
GOLD        = RGBColor(0xD4, 0xA8, 0x53)
CYAN        = RGBColor(0x00, 0xBC, 0xD4)
TEXT_MAIN   = RGBColor(0xF5, 0xF5, 0xF5)
TEXT_SUB    = RGBColor(0xB8, 0xB8, 0xC8)
CARD_BG     = RGBColor(0x1A, 0x1A, 0x2E)
TABLE_ROW   = RGBColor(0x12, 0x12, 0x24)
TABLE_ROW2  = RGBColor(0x16, 0x16, 0x28)

# ── Fonts ────────────────────────────────────────────────────────────────────
FONT_TITLE  = "Playfair Display"
FONT_BODY   = "Inter"

# ── Layout constants (inches) ────────────────────────────────────────────────
SLIDE_W     = 13.333
SLIDE_H     = 7.5
MARGIN_L    = 0.8
MARGIN_R    = 0.8
MARGIN_T    = 1.25
CONTENT_W   = SLIDE_W - MARGIN_L - MARGIN_R   # 11.733 in
BAR_H       = Inches(0.028)
FOOTER_Y    = Inches(7.05)
FOOTER_H    = Inches(0.35)
URL_Y       = Inches(7.1)

# ── EMU helpers ──────────────────────────────────────────────────────────────
def emu(inches): return Inches(inches)

# ── XML namespace ─────────────────────────────────────────────────────────────
nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
nsmap_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

def a(tag): return f"{{{nsmap_a}}}{tag}"
def p(tag): return f"{{{nsmap_p}}}{tag}"


# ── Presentation setup ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)

# Use blank layout (index 6 in default template)
blank_layout = prs.slide_layouts[6]


# ── Helper: set slide background ──────────────────────────────────────────────
def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── Helper: add thin horizontal bar ──────────────────────────────────────────
def add_bar(slide, y_inches, color=GOLD, width=None):
    w = width if width else CONTENT_W
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1 but we use freeform via add_shape
        emu(MARGIN_L), emu(y_inches), emu(w), emu(0.028)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


# ── Helper: add footer ────────────────────────────────────────────────────────
def add_footer(slide, url="theinnovativenative.com"):
    add_bar(slide, 7.05)
    tb = slide.shapes.add_textbox(
        emu(MARGIN_L), emu(7.1), emu(CONTENT_W), emu(0.35)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.RIGHT
    _set_para_defaults(p0, size=Pt(11), color=TEXT_SUB, font=FONT_BODY)
    run = p0.add_run()
    run.text = url


# ── Helper: paragraph XML defaults (size, font, color) ───────────────────────
def _set_para_defaults(para, size=None, color=None, font=None, bold=False, align=None):
    """Sets defRPr on a paragraph's pPr via XML manipulation."""
    pPr = para._pPr
    if pPr is None:
        pPr = etree.SubElement(para._p, a("pPr"))
    if align is not None:
        pPr.set("algn", {
            PP_ALIGN.LEFT: "l",
            PP_ALIGN.CENTER: "ctr",
            PP_ALIGN.RIGHT: "r",
        }.get(align, "l"))

    defRPr = pPr.find(a("defRPr"))
    if defRPr is None:
        defRPr = etree.SubElement(pPr, a("defRPr"))
    if size:
        defRPr.set("sz", str(int(size.pt * 100)))
    if bold:
        defRPr.set("b", "1")
    if color:
        solidFill = etree.SubElement(defRPr, a("solidFill"))
        srgbClr = etree.SubElement(solidFill, a("srgbClr"))
        srgbClr.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
    if font:
        latin = etree.SubElement(defRPr, a("latin"))
        latin.set("typeface", font)


def rgb_hex(color):
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


# ── Helper: add section header (gold bar + title textbox) ────────────────────
def add_section_header(slide, title_text, y_bar=1.25, y_title=1.35):
    add_bar(slide, y_bar)
    tb = slide.shapes.add_textbox(
        emu(MARGIN_L), emu(y_title), emu(CONTENT_W), emu(0.7)
    )
    tf = tb.text_frame
    tf.word_wrap = False
    para = tf.paragraphs[0]
    _set_para_defaults(para, size=Pt(28), color=GOLD, font=FONT_TITLE, bold=True)
    run = para.add_run()
    run.text = title_text


# ── Helper: build XML for a paragraph with run ────────────────────────────────
def xml_para(text, size_pt, color, font, bold=False, align="l", space_before=0, space_after=600):
    """Returns an <a:p> element."""
    p_el = etree.Element(a("p"))
    pPr = etree.SubElement(p_el, a("pPr"))
    pPr.set("algn", align)
    if space_before:
        spcBef = etree.SubElement(pPr, a("spcBef"))
        spcPts = etree.SubElement(spcBef, a("spcPts"))
        spcPts.set("val", str(space_before))
    spcAft = etree.SubElement(pPr, a("spcAft"))
    spcPtsAft = etree.SubElement(spcAft, a("spcPts"))
    spcPtsAft.set("val", str(space_after))
    defRPr = etree.SubElement(pPr, a("defRPr"))
    defRPr.set("sz", str(int(size_pt * 100)))
    if bold:
        defRPr.set("b", "1")
    sf = etree.SubElement(defRPr, a("solidFill"))
    cl = etree.SubElement(sf, a("srgbClr"))
    cl.set("val", rgb_hex(color))
    lat = etree.SubElement(defRPr, a("latin"))
    lat.set("typeface", font)

    r_el = etree.SubElement(p_el, a("r"))
    t_el = etree.SubElement(r_el, a("t"))
    t_el.text = text
    return p_el


def xml_para_mixed(parts, size_pt, base_color, font, align="l", space_before=0, space_after=600):
    """
    Returns <a:p> with mixed runs.
    parts = list of (text, color_override_or_None, bold_override_or_None)
    """
    p_el = etree.Element(a("p"))
    pPr = etree.SubElement(p_el, a("pPr"))
    pPr.set("algn", align)
    if space_before:
        spcBef = etree.SubElement(pPr, a("spcBef"))
        spcPts = etree.SubElement(spcBef, a("spcPts"))
        spcPts.set("val", str(space_before))
    spcAft = etree.SubElement(pPr, a("spcAft"))
    spcPtsAft = etree.SubElement(spcAft, a("spcPts"))
    spcPtsAft.set("val", str(space_after))
    defRPr = etree.SubElement(pPr, a("defRPr"))
    defRPr.set("sz", str(int(size_pt * 100)))
    sf = etree.SubElement(defRPr, a("solidFill"))
    cl = etree.SubElement(sf, a("srgbClr"))
    cl.set("val", rgb_hex(base_color))
    lat = etree.SubElement(defRPr, a("latin"))
    lat.set("typeface", font)

    for (text, color_ov, bold_ov) in parts:
        r_el = etree.SubElement(p_el, a("r"))
        rPr = etree.SubElement(r_el, a("rPr"))
        rPr.set("lang", "en-US")
        if bold_ov is not None:
            rPr.set("b", "1" if bold_ov else "0")
        if color_ov is not None:
            sf2 = etree.SubElement(rPr, a("solidFill"))
            cl2 = etree.SubElement(sf2, a("srgbClr"))
            cl2.set("val", rgb_hex(color_ov))
        t_el = etree.SubElement(r_el, a("t"))
        t_el.text = text
    return p_el


# ── Helper: add a table ───────────────────────────────────────────────────────
def add_table(slide, rows, cols, left, top, width, height, col_widths=None):
    """Add a table and return the table object. col_widths in inches."""
    tf = slide.shapes.add_table(rows, cols, emu(left), emu(top), emu(width), emu(height))
    tbl = tf.table

    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = emu(cw)

    return tbl


def style_cell(cell, text, fill_color, text_color, size_pt=14, bold=False,
               font=FONT_BODY, align=PP_ALIGN.LEFT, wrap=True):
    """Style a single table cell."""
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill_color

    tf = cell.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align

    # Build XML for the paragraph
    pPr = para._pPr
    if pPr is None:
        pPr = etree.SubElement(para._p, a("pPr"))
    pPr.set("algn", {PP_ALIGN.LEFT: "l", PP_ALIGN.CENTER: "ctr", PP_ALIGN.RIGHT: "r"}.get(align, "l"))

    defRPr = pPr.find(a("defRPr"))
    if defRPr is None:
        defRPr = etree.SubElement(pPr, a("defRPr"))
    defRPr.set("sz", str(int(size_pt * 100)))
    if bold:
        defRPr.set("b", "1")
    sf = etree.SubElement(defRPr, a("solidFill"))
    cl = etree.SubElement(sf, a("srgbClr"))
    cl.set("val", rgb_hex(text_color))
    lat = etree.SubElement(defRPr, a("latin"))
    lat.set("typeface", font)

    run = para.add_run()
    run.text = text

    # Cell margins
    tcPr = cell._tc.find(a("tcPr"))
    if tcPr is not None:
        tcPr.set("marL", "76200")
        tcPr.set("marR", "76200")
        tcPr.set("marT", "50800")
        tcPr.set("marB", "50800")


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title slide
# ─────────────────────────────────────────────────────────────────────────────
slide1 = prs.slides.add_slide(blank_layout)
set_bg(slide1)

# Center decorative bar
add_bar(slide1, 2.0)

# Main textbox (centered, vertically middle-ish)
tb = slide1.shapes.add_textbox(
    emu(MARGIN_L), emu(2.1), emu(CONTENT_W), emu(3.5)
)
tf = tb.text_frame
tf.word_wrap = True

# Title
p1 = tf.paragraphs[0]
p1.alignment = PP_ALIGN.CENTER
_set_para_defaults(p1, size=Pt(40), color=GOLD, font=FONT_TITLE, bold=True)
r1 = p1.add_run()
r1.text = "AI Demand Letter Tool"

p2 = tf.add_paragraph()
p2.alignment = PP_ALIGN.CENTER
_set_para_defaults(p2, size=Pt(40), color=GOLD, font=FONT_TITLE, bold=True)
r2 = p2.add_run()
r2.text = "Comparison Matrix + ROI Calculator"

# Subtitle
p3 = tf.add_paragraph()
p3.alignment = PP_ALIGN.CENTER
_set_para_defaults(p3, size=Pt(18), color=TEXT_SUB, font=FONT_BODY)
r3 = p3.add_run()
r3.text = "Side-by-side scoring of 4 tools. ROI math for your firm."

# Spacer paragraph
p4 = tf.add_paragraph()
_set_para_defaults(p4, size=Pt(14), color=TEXT_SUB, font=FONT_BODY, align=PP_ALIGN.CENTER)
r4 = p4.add_run()
r4.text = ""

# Attribution
p5 = tf.add_paragraph()
p5.alignment = PP_ALIGN.CENTER
_set_para_defaults(p5, size=Pt(14), color=TEXT_SUB, font=FONT_BODY)
r5 = p5.add_run()
r5.text = "Prepared by Mike Soto, The Innovative Native  |  April 2026"

p6 = tf.add_paragraph()
p6.alignment = PP_ALIGN.CENTER
_set_para_defaults(p6, size=Pt(13), color=TEXT_SUB, font=FONT_BODY)
r6 = p6.add_run()
r6.text = "For: PI firm owners, managing partners, operations leads"

add_footer(slide1)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — The Bottleneck Nobody Talks About
# ─────────────────────────────────────────────────────────────────────────────
slide2 = prs.slides.add_slide(blank_layout)
set_bg(slide2)

add_section_header(slide2, "The Bottleneck Nobody Talks About")

# Three stat callout boxes
callouts = [
    ("5–15 hrs",      "per demand letter"),
    ("200–600 hrs",   "per billing cycle"),
    ("$6K–$18K",      "in paralegal labor"),
]
box_w = 3.4
box_h = 1.4
box_gap = 0.27
box_start_x = MARGIN_L
box_y = 2.05

for i, (stat, label) in enumerate(callouts):
    bx = slide2.shapes.add_shape(
        1,
        emu(box_start_x + i * (box_w + box_gap)),
        emu(box_y),
        emu(box_w),
        emu(box_h)
    )
    bx.fill.solid()
    bx.fill.fore_color.rgb = CARD_BG
    bx.line.fill.background()

    # Stat number / amount
    stat_tb = slide2.shapes.add_textbox(
        emu(box_start_x + i * (box_w + box_gap) + 0.15),
        emu(box_y + 0.1),
        emu(box_w - 0.3),
        emu(0.7)
    )
    stf = stat_tb.text_frame
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    _set_para_defaults(sp, size=Pt(30), color=CYAN, font=FONT_TITLE, bold=True)
    sr = sp.add_run()
    sr.text = stat

    # Label
    lbl_tb = slide2.shapes.add_textbox(
        emu(box_start_x + i * (box_w + box_gap) + 0.15),
        emu(box_y + 0.82),
        emu(box_w - 0.3),
        emu(0.55)
    )
    ltf = lbl_tb.text_frame
    lp = ltf.paragraphs[0]
    lp.alignment = PP_ALIGN.CENTER
    _set_para_defaults(lp, size=Pt(13), color=TEXT_SUB, font=FONT_BODY)
    lr = lp.add_run()
    lr.text = label

# Body text
body_lines = [
    "Your paralegal spends 5 to 15 hours writing a single demand letter: gather records,",
    "summarize injuries, calculate damages, research settlements, draft, edit, review, send.",
    "",
    "Multiply that by 40 active cases: 200–600 hours per cycle. At $30/hr, that is",
    "$6,000–$18,000 in labor. Per round. This is not a talent problem — it's an automation problem.",
    "",
    "Harvard Law School (2025): tasks requiring 16 hours completed in 3–4 minutes with AI.",
    "A 100x productivity gain on your biggest operational bottleneck.",
    "",
    "The four tools below are shipping production demand letters in 2026. Pick one. Start this month.",
]

body_tb = slide2.shapes.add_textbox(
    emu(MARGIN_L), emu(3.65), emu(CONTENT_W), emu(3.1)
)
btf = body_tb.text_frame
btf.word_wrap = True

for li, line in enumerate(body_lines):
    para = btf.paragraphs[0] if li == 0 else btf.add_paragraph()
    _set_para_defaults(para, size=Pt(14), color=TEXT_MAIN if line else TEXT_SUB, font=FONT_BODY)
    run = para.add_run()
    run.text = line

add_footer(slide2)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — 4-Tool Matrix (table slide)
# ─────────────────────────────────────────────────────────────────────────────
slide3 = prs.slides.add_slide(blank_layout)
set_bg(slide3)

add_section_header(slide3, "4 Tools, Head to Head")

# Table data
headers = ["Tool", "Starting Cost", "Speed", "Best For", "Integrations"]
rows_data = [
    ["EvenUp",              "~$4M/yr enterprise;\ncustom for smaller",   "< 1 hr/letter",  "Mass tort +\nhigh-volume PI",        "Filevine, CloudLex,\nLitify"],
    ["Tavrn",               "$299/mo entry",                              "45 min/letter",  "Small-to-mid PI\n(10-25 attorneys)",  "Clio, CasePeer"],
    ["Filevine DemandsAI",  "Built into Filevine\n(~$75/user/mo)",        "1-2 hrs/letter", "Firms already\non Filevine",          "Native Filevine"],
    ["Supio",               "$1,500–$3,500/mo",                           "< 2 hrs/letter", "Mid-size firms\nwanting config",       "Most major CMS"],
]

col_widths = [2.2, 2.4, 1.7, 2.65, 2.73]
tbl = add_table(slide3, len(rows_data) + 1, 5,
                MARGIN_L, 2.0, CONTENT_W, 4.6, col_widths=col_widths)

# Header row
for ci, hdr in enumerate(headers):
    style_cell(tbl.cell(0, ci), hdr,
               fill_color=GOLD, text_color=RGBColor(0x0A, 0x0A, 0x14),
               size_pt=13, bold=True, font=FONT_BODY, align=PP_ALIGN.CENTER)

# Data rows
for ri, row in enumerate(rows_data):
    fill = TABLE_ROW if ri % 2 == 0 else TABLE_ROW2
    for ci, cell_text in enumerate(row):
        style_cell(tbl.cell(ri + 1, ci), cell_text,
                   fill_color=fill, text_color=TEXT_MAIN,
                   size_pt=12, font=FONT_BODY,
                   align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

# Headline numbers below table
hl_tb = slide3.shapes.add_textbox(
    emu(MARGIN_L), emu(6.65), emu(CONTENT_W), emu(0.35)
)
hltf = hl_tb.text_frame
hltf.word_wrap = False
hlp = hltf.paragraphs[0]
_set_para_defaults(hlp, size=Pt(11), color=TEXT_SUB, font=FONT_BODY)
hlr = hlp.add_run()
hlr.text = "EvenUp: 10,000 cases/week, $10B+ damages secured (Fortune 2026)  |  Harvard 2025: 16-hr task → 3-4 min with AI (100x)  |  Industry baseline: 5-15 hrs/letter (ProPlaintiff 2026)"

add_footer(slide3)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — ROI Calculator
# ─────────────────────────────────────────────────────────────────────────────
slide4 = prs.slides.add_slide(blank_layout)
set_bg(slide4)

add_section_header(slide4, "ROI Calculator: Your Firm's Real Numbers")

# Three-step framework cards
steps = [
    ("Step 1",  "BASELINE",  "Cases/mo × Hrs/letter × $/hr\n= Monthly manual cost"),
    ("Step 2",  "WITH AI",   "Tool cost + Review hrs × $/hr\n= Monthly AI cost"),
    ("Step 3",  "SAVINGS",   "Manual cost − AI cost\n= Monthly savings"),
]
card_w  = 3.5
card_h  = 1.8
card_gap= 0.24
card_y  = 2.05

for i, (step_lbl, step_title, step_body) in enumerate(steps):
    cx = MARGIN_L + i * (card_w + card_gap)

    card = slide4.shapes.add_shape(
        1,
        emu(cx), emu(card_y), emu(card_w), emu(card_h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.fill.background()

    # Step label
    step_tb = slide4.shapes.add_textbox(
        emu(cx + 0.12), emu(card_y + 0.1), emu(card_w - 0.24), emu(0.35)
    )
    stp = step_tb.text_frame.paragraphs[0]
    stp.alignment = PP_ALIGN.CENTER
    _set_para_defaults(stp, size=Pt(11), color=GOLD, font=FONT_BODY, bold=True)
    stp.add_run().text = step_lbl

    # Step title
    title_tb = slide4.shapes.add_textbox(
        emu(cx + 0.12), emu(card_y + 0.42), emu(card_w - 0.24), emu(0.45)
    )
    ttp = title_tb.text_frame.paragraphs[0]
    ttp.alignment = PP_ALIGN.CENTER
    _set_para_defaults(ttp, size=Pt(18), color=CYAN, font=FONT_TITLE, bold=True)
    ttp.add_run().text = step_title

    # Step body
    body_tb2 = slide4.shapes.add_textbox(
        emu(cx + 0.12), emu(card_y + 0.95), emu(card_w - 0.24), emu(0.75)
    )
    body_tb2.text_frame.word_wrap = True
    bp = body_tb2.text_frame.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    _set_para_defaults(bp, size=Pt(12), color=TEXT_SUB, font=FONT_BODY)
    bp.add_run().text = step_body

# Arrow connectors (simple text "→")
for i in range(2):
    arr_tb = slide4.shapes.add_textbox(
        emu(MARGIN_L + (i + 1) * (card_w + card_gap) - card_gap + 0.04),
        emu(card_y + card_h / 2 - 0.2),
        emu(0.2),
        emu(0.4)
    )
    ap = arr_tb.text_frame.paragraphs[0]
    ap.alignment = PP_ALIGN.CENTER
    _set_para_defaults(ap, size=Pt(18), color=GOLD, font=FONT_BODY, bold=True)
    ap.add_run().text = "→"

# Worked example box
ex_box = slide4.shapes.add_shape(
    1,
    emu(MARGIN_L), emu(4.05), emu(CONTENT_W), emu(2.4)
)
ex_box.fill.solid()
ex_box.fill.fore_color.rgb = CARD_BG
ex_box.line.fill.background()

# Gold label bar for example
ex_label_bar = slide4.shapes.add_shape(
    1,
    emu(MARGIN_L), emu(4.05), emu(CONTENT_W), emu(0.03)
)
ex_label_bar.fill.solid()
ex_label_bar.fill.fore_color.rgb = GOLD
ex_label_bar.line.fill.background()

ex_tb = slide4.shapes.add_textbox(
    emu(MARGIN_L + 0.3), emu(4.1), emu(CONTENT_W - 0.6), emu(2.3)
)
ex_tf = ex_tb.text_frame
ex_tf.word_wrap = True

ex_lines = [
    ("WORKED EXAMPLE — 25-attorney PI firm", GOLD, True, Pt(13)),
    ("40 demand letters/mo · $30/hr paralegal · 10 hrs per letter", TEXT_SUB, False, Pt(12)),
    ("", TEXT_SUB, False, Pt(8)),
    ("Manual:    40 × 10 × $30 = $12,000/mo  ($144,000/yr)", TEXT_MAIN, False, Pt(13)),
    ("With Tavrn: 40 × 1.5 × $30 + $299 = $2,099/mo  ($25,000/yr)", TEXT_MAIN, False, Pt(13)),
    ("Savings:   $9,901/mo · $119,000/yr · Break-even: < 1 letter", CYAN, True, Pt(14)),
    ("", TEXT_SUB, False, Pt(8)),
    ('"The easiest business-case math most PI firms will ever run."', TEXT_SUB, False, Pt(12)),
]

for li, (text, color, bold, sz) in enumerate(ex_lines):
    ep = ex_tf.paragraphs[0] if li == 0 else ex_tf.add_paragraph()
    ep.alignment = PP_ALIGN.LEFT
    _set_para_defaults(ep, size=sz, color=color, font=FONT_BODY, bold=bold)
    ep.add_run().text = text

add_footer(slide4)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — The 5 Red Flags
# ─────────────────────────────────────────────────────────────────────────────
slide5 = prs.slides.add_slide(blank_layout)
set_bg(slide5)

add_section_header(slide5, "5 Red Flags in Every Demo")

sub_tb = slide5.shapes.add_textbox(
    emu(MARGIN_L), emu(2.0), emu(CONTENT_W), emu(0.35)
)
sp = sub_tb.text_frame.paragraphs[0]
_set_para_defaults(sp, size=Pt(14), color=TEXT_SUB, font=FONT_BODY)
sp.add_run().text = "When you sit through a sales pitch — if the vendor dodges any of these, walk."

flags = [
    (
        "1",
        '"Can you show me a live generation on a case I brought?"',
        "If they want to use their canned demo, they are hiding output quality. A real tool runs on real records in real time."
    ),
    (
        "2",
        '"Who owns the data after it enters your system?"',
        'If the answer is "we retain it for training," your privileged info just left your firm. Hard no after Heppner.'
    ),
    (
        "3",
        '"What is your hallucination rate on citations?"',
        'Published benchmarks or walk. "Very low" is not a number.'
    ),
    (
        "4",
        '"Show me the paralegal review workflow."',
        "The right answer: AI drafts, paralegal verifies, attorney signs. Anything else is pre-Heppner thinking."
    ),
    (
        "5",
        '"What integrations do you have with my CMS?"',
        "Filevine, Clio, CasePeer, SmartAdvocate, CloudLex. If your CMS isn't on their list, the rollout will sink."
    ),
]

flag_y = 2.45
flag_h = 0.82
flag_gap = 0.06

for num, question, explanation in flags:
    # Number badge
    badge = slide5.shapes.add_shape(
        1,
        emu(MARGIN_L), emu(flag_y), emu(0.42), emu(flag_h)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = GOLD
    badge.line.fill.background()

    num_tb = slide5.shapes.add_textbox(
        emu(MARGIN_L), emu(flag_y + 0.15), emu(0.42), emu(0.42)
    )
    np = num_tb.text_frame.paragraphs[0]
    np.alignment = PP_ALIGN.CENTER
    _set_para_defaults(np, size=Pt(18), color=RGBColor(0x0A, 0x0A, 0x14), font=FONT_TITLE, bold=True)
    np.add_run().text = num

    # Content textbox
    content_tb = slide5.shapes.add_textbox(
        emu(MARGIN_L + 0.55), emu(flag_y + 0.04), emu(CONTENT_W - 0.55), emu(flag_h - 0.04)
    )
    ctf = content_tb.text_frame
    ctf.word_wrap = True

    qp = ctf.paragraphs[0]
    _set_para_defaults(qp, size=Pt(13), color=GOLD, font=FONT_BODY, bold=True)
    qp.add_run().text = question

    ep = ctf.add_paragraph()
    _set_para_defaults(ep, size=Pt(12), color=TEXT_SUB, font=FONT_BODY)
    ep.add_run().text = explanation

    flag_y += flag_h + flag_gap

add_footer(slide5)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Monday Morning Action List
# ─────────────────────────────────────────────────────────────────────────────
slide6 = prs.slides.add_slide(blank_layout)
set_bg(slide6)

add_section_header(slide6, "Monday Morning. 3 Actions.")

sub2_tb = slide6.shapes.add_textbox(
    emu(MARGIN_L), emu(2.0), emu(CONTENT_W), emu(0.35)
)
s2p = sub2_tb.text_frame.paragraphs[0]
_set_para_defaults(s2p, size=Pt(14), color=TEXT_SUB, font=FONT_BODY)
s2p.add_run().text = "Three things you can do before lunch on the first work day after reading this:"

actions = [
    (
        "1",
        "Time your current process.",
        "Email your paralegal. Ask them to track hours on the next three demand letters. "
        "Don't tell them why. Get the real number, not the estimate."
    ),
    (
        "2",
        "Pull your last 10 closed case files.",
        "Count how many demand letters were involved. Calculate actual labor hours at your paralegal rate. "
        "That is your monthly demand letter cost, minimum."
    ),
    (
        "3",
        "Book a 30-minute Tavrn demo.",
        "$299/mo is cheap enough that you don't need your CFO to sign off. "
        "Walk through one live generation on a de-identified case file. Decide by Friday."
    ),
]

action_y = 2.5
action_h = 1.3
action_gap = 0.18

for num, title, body in actions:
    # Number badge
    badge2 = slide6.shapes.add_shape(
        1,
        emu(MARGIN_L), emu(action_y), emu(0.5), emu(action_h)
    )
    badge2.fill.solid()
    badge2.fill.fore_color.rgb = GOLD
    badge2.line.fill.background()

    num_tb2 = slide6.shapes.add_textbox(
        emu(MARGIN_L), emu(action_y + 0.3), emu(0.5), emu(0.5)
    )
    np2 = num_tb2.text_frame.paragraphs[0]
    np2.alignment = PP_ALIGN.CENTER
    _set_para_defaults(np2, size=Pt(22), color=RGBColor(0x0A, 0x0A, 0x14), font=FONT_TITLE, bold=True)
    np2.add_run().text = num

    # Content textbox
    act_tb = slide6.shapes.add_textbox(
        emu(MARGIN_L + 0.65), emu(action_y + 0.1), emu(CONTENT_W - 0.65), emu(action_h - 0.1)
    )
    atf = act_tb.text_frame
    atf.word_wrap = True

    atp = atf.paragraphs[0]
    _set_para_defaults(atp, size=Pt(15), color=GOLD, font=FONT_BODY, bold=True)
    atp.add_run().text = title

    abp = atf.add_paragraph()
    _set_para_defaults(abp, size=Pt(13), color=TEXT_MAIN, font=FONT_BODY)
    abp.add_run().text = body

    action_y += action_h + action_gap

# CTA
cta_box = slide6.shapes.add_shape(
    1,
    emu(MARGIN_L), emu(6.45), emu(CONTENT_W), emu(0.45)
)
cta_box.fill.solid()
cta_box.fill.fore_color.rgb = CARD_BG
cta_box.line.fill.background()

cta_tb = slide6.shapes.add_textbox(
    emu(MARGIN_L + 0.2), emu(6.47), emu(CONTENT_W - 0.4), emu(0.4)
)
ctf2 = cta_tb.text_frame
ctf2.word_wrap = False
cp = ctf2.paragraphs[0]
cp.alignment = PP_ALIGN.CENTER
_set_para_defaults(cp, size=Pt(13), color=CYAN, font=FONT_BODY, bold=True)
cp.add_run().text = "buildmytribe.ai/demand  |  theinnovativenative.com/resources/demand-letter-matrix"

# Sources line
src_tb = slide6.shapes.add_textbox(
    emu(MARGIN_L), emu(6.95), emu(CONTENT_W), emu(0.3)
)
sp3 = src_tb.text_frame.paragraphs[0]
sp3.alignment = PP_ALIGN.LEFT
_set_para_defaults(sp3, size=Pt(10), color=TEXT_SUB, font=FONT_BODY)
sp3.add_run().text = (
    "Sources: EvenUp Series E (Fortune 2026)  |  Harvard Law / MIT Sloan 2025  |  "
    "ProPlaintiff 2026  |  Tavrn, Filevine, Supio product specs 2026  |  Mike Soto, The Innovative Native"
)

add_footer(slide6, url="theinnovativenative.com")


# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
output_path = (
    "/Users/makwa/theinnovativenative/projects/website/public/assets/lead-magnets/"
    "demand-letter-matrix.pptx"
)
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")

import os
size_kb = os.path.getsize(output_path) // 1024
print(f"Size: {size_kb} KB")

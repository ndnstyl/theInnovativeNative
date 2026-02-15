#!/usr/bin/env python3
"""
BATCH 1: Slides 1-5
- Slide 1: Title Slide
- Slide 2: Content Slide (What This Is/Is Not)
- Slide 3: Section Break (The $100K Problem)
- Slide 4: Stats Slide (Revenue Loss Breakdown)
- Slide 5: Two-Column Slide (Before vs After)
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGB"""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

# Brand colors from The Innovative Native
BRAND_BG = "0a0a14"
BRAND_BG_ALT = "12121f"
BRAND_TEXT = "f5f5f5"
BRAND_TEXT_SECONDARY = "b8b8c8"
BRAND_ACCENT = "00ffff"
BRAND_ACCENT_SECONDARY = "ff00ff"
BRAND_ACCENT_TERTIARY = "00d4d4"
BRAND_CARD_BG = "1a1a2e"
BRAND_CARD_BG_ALT = "16213e"
BRAND_HEADING_FONT = "Playfair Display"
BRAND_BODY_FONT = "Inter"

def create_batch1():
    """Generate slides 1-5"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 1: Title Slide =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Accent bar at top
    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.15)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    top_bar.line.fill.background()

    # Headline
    headline_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(2.8),
        Inches(11.83), Inches(1.8)
    )
    tf = headline_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The Automated Law Firm Blueprint"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.8),
        Inches(11.83), Inches(1.2)
    )
    p = subtitle_box.text_frame.paragraphs[0]
    p.text = "A Strategic Architecture for Bankruptcy, Criminal & Administrative Law Firms"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
    p.alignment = PP_ALIGN.CENTER

    # ===== SLIDE 2: Content Slide - What This Is/Is Not =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Left accent bar
    left_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(0.08), Inches(7.5)
    )
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    left_bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.5),
        Inches(12.0), Inches(1.0)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = "Know Before You Read"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Accent underline
    underline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.45),
        Inches(2.5), Inches(0.05)
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    underline.line.fill.background()

    # Content bullets
    content_box = slide.shapes.add_textbox(
        Inches(0.6), Inches(2.2),
        Inches(12.0), Inches(4.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    bullets = [
        "What This Is: Strategic document showing a fully automated law firm — from intake to marketing",
        "What This Is Not: A course, tutorial, or legal advice. We build systems for law firms",
        "Scope: Practice-area architecture for Bankruptcy, Criminal, and Administrative law",
        "Disclaimer: We build operational and marketing infrastructure — the business side of running a firm"
    ]

    for i, bullet in enumerate(bullets):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = bullet
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(16)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
        p.space_before = Pt(10) if i > 0 else Pt(0)
        p.level = 0

    # ===== SLIDE 3: Section Break - The $100K Problem =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Large accent block
    accent_block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(4.5), Inches(7.5)
    )
    accent_block.fill.solid()
    accent_block.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT)
    accent_block.line.fill.background()

    # Section number
    num_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.8),
        Inches(3.5), Inches(1.5)
    )
    p = num_box.text_frame.paragraphs[0]
    p.text = "01"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
    p.alignment = PP_ALIGN.CENTER

    # Section title
    title_box = slide.shapes.add_textbox(
        Inches(5.0), Inches(3.0),
        Inches(7.5), Inches(1.5)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = "The $100K Problem"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # ===== SLIDE 4: Stats Slide - Revenue Loss Breakdown =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Top accent line
    top_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.3),
        Inches(3.0), Inches(0.06)
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    top_line.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.6),
        Inches(12.0), Inches(1.0)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = "Where Your Revenue Goes Every Year"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Stats
    stats = [
        ("35%", "Missed Billable Time"),
        ("25%", "Slow Intake Conversion"),
        ("20%", "Manual Document Work"),
        ("15%", "Ineffective Marketing"),
        ("5%", "Admin Overhead")
    ]

    start_y = 2.2
    spacing = 1.0
    for i, (value, label) in enumerate(stats):
        y_pos = start_y + (i * spacing)

        # Value
        value_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(y_pos),
            Inches(3.0), Inches(0.7)
        )
        p = value_box.text_frame.paragraphs[0]
        p.text = value
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(48)
        p.font.bold = True
        accent_color = BRAND_ACCENT if i < 2 else (BRAND_ACCENT_SECONDARY if i < 4 else BRAND_TEXT_SECONDARY)
        p.font.color.rgb = hex_to_rgb(accent_color)

        # Label
        label_box = slide.shapes.add_textbox(
            Inches(4.8), Inches(y_pos + 0.15),
            Inches(6.5), Inches(0.5)
        )
        p = label_box.text_frame.paragraphs[0]
        p.text = label
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # ===== SLIDE 5: Two-Column Slide - Before vs After =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4),
        Inches(12.333), Inches(0.9)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = "The Transformation"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Left column header
    left_header_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.6),
        Inches(5.5), Inches(0.6)
    )
    p = left_header_box.text_frame.paragraphs[0]
    p.text = "Before Automation"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # Left bullets
    left_bullets = [
        "Sticky notes, manual follow-ups",
        "Billing from memory at end of week",
        "3 hours drafting repetitive motions",
        "Clients calling: 'Where's my case?'"
    ]
    left_content = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.4),
        Inches(5.5), Inches(4.0)
    )
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = bullet
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(16)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
        p.space_before = Pt(12) if i > 0 else Pt(0)
        p.level = 0

    # Right column header
    right_header_box = slide.shapes.add_textbox(
        Inches(7.0), Inches(1.6),
        Inches(5.5), Inches(0.6)
    )
    p = right_header_box.text_frame.paragraphs[0]
    p.text = "After Automation"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    # Right bullets
    right_bullets = [
        "Instant conflict check, auto-response",
        "Activity-based capture: 95% recovery",
        "45 minutes review, templates handle rest",
        "Automated status updates at milestones"
    ]
    right_content = slide.shapes.add_textbox(
        Inches(7.0), Inches(2.4),
        Inches(5.5), Inches(4.0)
    )
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i > 0:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = bullet
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(16)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
        p.space_before = Pt(12) if i > 0 else Pt(0)
        p.level = 0

    return prs

if __name__ == "__main__":
    prs = create_batch1()
    prs.save('part1.pptx')
    print("✓ Batch 1 complete: Slides 1-5 saved to part1.pptx")

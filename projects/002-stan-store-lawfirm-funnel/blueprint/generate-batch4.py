#!/usr/bin/env python3
"""
BATCH 4: Slides 16-20
- Slide 16: Multi-Card (Communication automation touchpoints)
- Slide 17: Quote (Client communication insight)
- Slide 18: Section Break (The Billing Machine)
- Slide 19: Stats (Billing recovery metrics)
- Slide 20: Two-Column (Manual vs automated billing)
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

BRAND_BG = "0a0a14"
BRAND_BG_ALT = "12121f"
BRAND_TEXT = "f5f5f5"
BRAND_TEXT_SECONDARY = "b8b8c8"
BRAND_ACCENT = "00ffff"
BRAND_ACCENT_SECONDARY = "ff00ff"
BRAND_CARD_BG = "1a1a2e"
BRAND_HEADING_FONT = "Playfair Display"
BRAND_BODY_FONT = "Inter"

def create_batch4():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 16: Multi-Card - Communication Touchpoints =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Automated Communication Touchpoints"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    touchpoints = [
        ("Status Updates", "Automatic notifications at\nevery case milestone —\nfiling, hearing, ruling"),
        ("Appointment Reminders", "48-hour and 2-hour\nreminders with prep\ninstructions attached"),
        ("Document Requests", "Automated follow-ups\nfor missing documents\nwith secure upload links"),
        ("Check-Ins", "Periodic case updates\neven when nothing changes —\nsilence breeds anxiety"),
    ]
    for i, (title, desc) in enumerate(touchpoints):
        x = 0.5 + (i * 3.2)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(2.8), Inches(4.0))
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
        card.line.color.rgb = hex_to_rgb(BRAND_ACCENT)
        card.line.width = Pt(1)

        t_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(2.5), Inches(2.2), Inches(0.6))
        p = t_box.text_frame.paragraphs[0]
        p.text = title
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

        d_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(3.3), Inches(2.2), Inches(2.2))
        d_box.text_frame.word_wrap = True
        p = d_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # ===== SLIDE 17: Quote Slide =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Large quotation mark
    q_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.0), Inches(2.0), Inches(2.0))
    p = q_box.text_frame.paragraphs[0]
    p.text = "\u201C"
    p.font.size = Pt(200)
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)
    p.font.name = BRAND_HEADING_FONT

    # Quote text
    quote_box = slide.shapes.add_textbox(Inches(2.0), Inches(2.5), Inches(9.0), Inches(2.5))
    quote_box.text_frame.word_wrap = True
    p = quote_box.text_frame.paragraphs[0]
    p.text = "The number one driver of bar complaints is not incompetence. It is lack of communication. Clients do not complain because their lawyer lost. They complain because their lawyer disappeared."
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(28)
    p.font.italic = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Attribution
    attr_box = slide.shapes.add_textbox(Inches(2.0), Inches(5.3), Inches(9.0), Inches(0.5))
    p = attr_box.text_frame.paragraphs[0]
    p.text = "— ABA Research on Client Satisfaction"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(16)
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    # ===== SLIDE 18: Section Break - The Billing Machine =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    num_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(4.0), Inches(2.0))
    p = num_box.text_frame.paragraphs[0]
    p.text = "06"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.8), Inches(3.0), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(10.0), Inches(1.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = "The Billing Machine"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(10.0), Inches(0.8))
    p = sub_box.text_frame.paragraphs[0]
    p.text = "Stop billing from memory — capture every minute automatically"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # ===== SLIDE 19: Stats - Billing Recovery Metrics =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = "The Billing Recovery Opportunity"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    stats = [
        ("$109K", "Lost annually to untracked billable time at $350/hr"),
        ("15-30%", "Of billable hours lost when billing from memory"),
        ("95%", "Time capture rate with activity-based automation"),
    ]
    for i, (value, label) in enumerate(stats):
        y = 2.0 + (i * 1.6)
        v_box = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(3.5), Inches(0.7))
        p = v_box.text_frame.paragraphs[0]
        p.text = value
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(52)
        p.font.bold = True
        colors = [BRAND_ACCENT, BRAND_ACCENT_SECONDARY, BRAND_ACCENT]
        p.font.color.rgb = hex_to_rgb(colors[i])

        l_box = slide.shapes.add_textbox(Inches(5.5), Inches(y + 0.1), Inches(7.0), Inches(0.5))
        l_box.text_frame.word_wrap = True
        p = l_box.text_frame.paragraphs[0]
        p.text = label
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # ===== SLIDE 20: Two-Column - Manual vs Automated Billing =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Billing: Before and After"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Left header
    lh_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.6))
    p = lh_box.text_frame.paragraphs[0]
    p.text = "Manual Billing"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    left_items = ["Bill from memory end-of-week", "Lose 15-30% of billable hours", "Invoice generation takes hours", "Collections require manual follow-up", "No visibility into realization rates"]
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(5.5), Inches(4.0))
    left_box.text_frame.word_wrap = True
    for i, item in enumerate(left_items):
        p = left_box.text_frame.paragraphs[0] if i == 0 else left_box.text_frame.add_paragraph()
        p.text = item
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(16)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
        p.space_before = Pt(12) if i > 0 else Pt(0)

    # Right header
    rh_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.6))
    p = rh_box.text_frame.paragraphs[0]
    p.text = "Automated Billing"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    right_items = ["Activity-based time capture", "95% billable hour recovery", "One-click invoice generation", "Automated payment reminders", "Real-time realization dashboards"]
    right_box = slide.shapes.add_textbox(Inches(7.0), Inches(2.4), Inches(5.5), Inches(4.0))
    right_box.text_frame.word_wrap = True
    for i, item in enumerate(right_items):
        p = right_box.text_frame.paragraphs[0] if i == 0 else right_box.text_frame.add_paragraph()
        p.text = item
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(16)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
        p.space_before = Pt(12) if i > 0 else Pt(0)

    return prs

if __name__ == "__main__":
    prs = create_batch4()
    prs.save("part4.pptx")
    print("✓ Batch 4 complete: Slides 16-20 saved to part4.pptx")

#!/usr/bin/env python3
"""
BATCH 2: Slides 6-10
- Slide 6: Section Break (Automated Intake Engine)
- Slide 7: Floating-Cards (3 intake problems)
- Slide 8: Multi-Card (Intake automation architecture)
- Slide 9: Two-Column (Practice-area intake differences)
- Slide 10: Section Break (Document Assembly System)
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

BRAND_BG = "0a0a14"
BRAND_BG_ALT = "12121f"
BRAND_TEXT = "f5f5f5"
BRAND_TEXT_SECONDARY = "b8b8c8"
BRAND_ACCENT = "00ffff"
BRAND_ACCENT_SECONDARY = "ff00ff"
BRAND_CARD_BG = "1a1a2e"
BRAND_HEADING_FONT = "Playfair Display"
BRAND_BODY_FONT = "Inter"

def create_batch2():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 6: Section Break - Automated Intake Engine =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    # Section number
    num_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(4.0), Inches(2.0))
    p = num_box.text_frame.paragraphs[0]
    p.text = "02"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.8), Inches(3.0), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    line.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(10.0), Inches(1.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = "The Automated Intake Engine"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(10.0), Inches(0.8))
    p = sub_box.text_frame.paragraphs[0]
    p.text = "From sticky notes to sub-5-minute response times"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # ===== SLIDE 7: Floating-Cards - 3 Intake Problems =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Three Ways Intake Kills Your Pipeline"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    cards = [
        ("Speed-to-Lead", "Firms that respond within 5 minutes are 21x more likely to convert. Most firms respond in hours or days."),
        ("Conflict Checks", "Manual conflict checks take 15-30 minutes. Automated checks take seconds and never miss a match."),
        ("Scheduling Friction", "Every email exchange to schedule a consultation is a chance for the client to call someone else."),
    ]
    card_positions = [(0.5, 1.8), (4.6, 2.2), (8.7, 1.8)]
    for i, ((title, desc), (x, y)) in enumerate(zip(cards, card_positions)):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.0), Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
        card.line.color.rgb = hex_to_rgb(BRAND_ACCENT if i == 0 else BRAND_ACCENT_SECONDARY if i == 1 else "00d4d4")
        card.line.width = Pt(1.5)

        t_box = slide.shapes.add_textbox(Inches(x + 0.4), Inches(y + 0.5), Inches(3.2), Inches(0.6))
        p = t_box.text_frame.paragraphs[0]
        p.text = title
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

        d_box = slide.shapes.add_textbox(Inches(x + 0.4), Inches(y + 1.4), Inches(3.2), Inches(2.5))
        d_box.text_frame.word_wrap = True
        p = d_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(15)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # ===== SLIDE 8: Multi-Card - Intake Automation Architecture =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = "The Intake Engine Architecture"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    steps = [
        ("Web Form", "Practice-area specific,\nmobile-responsive, 24/7"),
        ("Conflict Check", "Instant search against\nentire client database"),
        ("E-Sign + Payment", "Auto-generated engagement\nletter + online retainer"),
        ("Case File Creation", "All intake data\npre-populated automatically"),
    ]
    for i, (title, desc) in enumerate(steps):
        x = 0.5 + (i * 3.2)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(2.8), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
        card.line.color.rgb = hex_to_rgb(BRAND_ACCENT)
        card.line.width = Pt(1)

        # Step number
        n_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(2.3), Inches(1.0), Inches(0.6))
        p = n_box.text_frame.paragraphs[0]
        p.text = str(i + 1)
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

        t_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(3.0), Inches(2.2), Inches(0.5))
        p = t_box.text_frame.paragraphs[0]
        p.text = title
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

        d_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(3.6), Inches(2.2), Inches(1.5))
        d_box.text_frame.word_wrap = True
        p = d_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

        # Arrow between cards
        if i < 3:
            arrow_box = slide.shapes.add_textbox(Inches(x + 2.9), Inches(3.3), Inches(0.4), Inches(0.5))
            p = arrow_box.text_frame.paragraphs[0]
            p.text = "→"
            p.font.size = Pt(28)
            p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    # ===== SLIDE 9: Two-Column - Practice-Area Intake Differences =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Intake by Practice Area"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Three columns for three practice areas
    areas = [
        ("Bankruptcy", ["Means test pre-screening", "Creditor list upload", "District-specific forms", "Petition data collection"]),
        ("Criminal Defense", ["Charge assessment", "Jurisdiction routing", "Bail/bond status", "Co-defendant conflict check"]),
        ("Administrative", ["Agency identification", "Deadline calculator", "Standing verification", "Document request checklist"]),
    ]
    for i, (area, items) in enumerate(areas):
        x = 0.5 + (i * 4.2)
        # Header
        h_box = slide.shapes.add_textbox(Inches(x), Inches(1.6), Inches(3.8), Inches(0.6))
        p = h_box.text_frame.paragraphs[0]
        p.text = area
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(22)
        p.font.bold = True
        colors = [BRAND_ACCENT, BRAND_ACCENT_SECONDARY, "00d4d4"]
        p.font.color.rgb = hex_to_rgb(colors[i])

        # Items
        items_box = slide.shapes.add_textbox(Inches(x), Inches(2.5), Inches(3.8), Inches(4.0))
        items_box.text_frame.word_wrap = True
        for j, item in enumerate(items):
            if j > 0:
                p = items_box.text_frame.add_paragraph()
            else:
                p = items_box.text_frame.paragraphs[0]
            p.text = f"→ {item}"
            p.font.name = BRAND_BODY_FONT
            p.font.size = Pt(16)
            p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)
            p.space_before = Pt(12) if j > 0 else Pt(0)

    # ===== SLIDE 10: Section Break - Document Assembly System =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    num_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(4.0), Inches(2.0))
    p = num_box.text_frame.paragraphs[0]
    p.text = "03"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.8), Inches(3.0), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(10.0), Inches(1.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Document Assembly System"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(10.0), Inches(0.8))
    p = sub_box.text_frame.paragraphs[0]
    p.text = "80% of your documents are templates waiting to happen"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    return prs

if __name__ == "__main__":
    prs = create_batch2()
    prs.save("part2.pptx")
    print("✓ Batch 2 complete: Slides 6-10 saved to part2.pptx")

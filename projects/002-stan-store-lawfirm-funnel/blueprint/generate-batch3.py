#!/usr/bin/env python3
"""
BATCH 3: Slides 11-15
- Slide 11: Multi-Card (Document automation components)
- Slide 12: Stats (Time savings per practice area)
- Slide 13: Section Break (Deadline & Compliance Engine)
- Slide 14: Floating-Cards (Compliance requirements)
- Slide 15: Section Break (Client Communication Hub)
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

BRAND_BG = "0a0a14"
BRAND_BG_ALT = "12121f"
BRAND_TEXT = "f5f5f5"
BRAND_TEXT_SECONDARY = "b8b8c8"
BRAND_ACCENT = "00ffff"
BRAND_ACCENT_SECONDARY = "ff00ff"
BRAND_CARD_BG = "1a1a2e"
BRAND_HEADING_FONT = "Playfair Display"
BRAND_BODY_FONT = "Inter"

def create_batch3():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 11: Multi-Card - Document Automation Components =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Document Assembly Components"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    components = [
        ("Template Library", "Standardized templates for\nevery document type in\nyour practice area"),
        ("Merge Fields", "Client data auto-populates\nnames, dates, case numbers,\njurisdiction details"),
        ("Assembly Engine", "Conditional logic selects\ncorrect clauses based on\ncase type and jurisdiction"),
        ("Quality Check", "Automated review catches\nmissing fields, formatting\nerrors, citation issues"),
    ]
    for i, (title, desc) in enumerate(components):
        x = 0.5 + (i * 3.2)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(2.8), Inches(4.0))
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
        card.line.color.rgb = hex_to_rgb(BRAND_ACCENT)
        card.line.width = Pt(1)

        t_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(2.5), Inches(2.2), Inches(0.6))
        p = t_box.text_frame.paragraphs[0]
        p.text = title
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

        d_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(3.3), Inches(2.2), Inches(2.2))
        d_box.text_frame.word_wrap = True
        p = d_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # ===== SLIDE 12: Stats - Time Savings Per Practice Area =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Document Assembly Saves Hours Every Week"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    stats = [
        ("7 hrs/wk", "Bankruptcy", "Petitions, schedules, means test forms — 80% template-driven"),
        ("5 hrs/wk", "Criminal Defense", "Motions to suppress, discovery requests, plea agreements"),
        ("6 hrs/wk", "Administrative", "Agency responses, FOIA requests, regulatory filings"),
    ]
    for i, (value, label, desc) in enumerate(stats):
        y = 2.0 + (i * 1.7)
        v_box = slide.shapes.add_textbox(Inches(1.0), Inches(y), Inches(3.0), Inches(0.7))
        p = v_box.text_frame.paragraphs[0]
        p.text = value
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

        l_box = slide.shapes.add_textbox(Inches(4.5), Inches(y), Inches(3.0), Inches(0.5))
        p = l_box.text_frame.paragraphs[0]
        p.text = label
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

        d_box = slide.shapes.add_textbox(Inches(4.5), Inches(y + 0.5), Inches(7.5), Inches(0.5))
        d_box.text_frame.word_wrap = True
        p = d_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(14)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # ===== SLIDE 13: Section Break - Deadline & Compliance Engine =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    num_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(4.0), Inches(2.0))
    p = num_box.text_frame.paragraphs[0]
    p.text = "04"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.8), Inches(3.0), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(10.0), Inches(1.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Deadline & Compliance Engine"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(10.0), Inches(0.8))
    p = sub_box.text_frame.paragraphs[0]
    p.text = "Missing a deadline doesn't just lose a case — it can end a career"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # ===== SLIDE 14: Floating-Cards - Compliance by Practice Area =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Critical Compliance by Practice Area"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    cards = [
        ("Bankruptcy", "Filing deadlines, credit counseling certificates, means test updates, trustee meeting schedules, plan payment monitoring"),
        ("Criminal Defense", "Speedy trial clocks, motion filing windows, discovery deadlines, plea agreement timelines, sentencing guidelines"),
        ("Administrative", "Comment period deadlines, appeal windows, hearing dates, agency response requirements, FOIA timelines"),
    ]
    card_positions = [(0.5, 1.8), (4.6, 2.2), (8.7, 1.8)]
    colors = [BRAND_ACCENT, BRAND_ACCENT_SECONDARY, "00d4d4"]
    for i, ((title, desc), (x, y)) in enumerate(zip(cards, card_positions)):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.0), Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG)
        card.line.color.rgb = hex_to_rgb(colors[i])
        card.line.width = Pt(1.5)

        t_box = slide.shapes.add_textbox(Inches(x + 0.4), Inches(y + 0.5), Inches(3.2), Inches(0.6))
        p = t_box.text_frame.paragraphs[0]
        p.text = title
        p.font.name = BRAND_HEADING_FONT
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(colors[i])

        d_box = slide.shapes.add_textbox(Inches(x + 0.4), Inches(y + 1.4), Inches(3.2), Inches(2.5))
        d_box.text_frame.word_wrap = True
        p = d_box.text_frame.paragraphs[0]
        p.text = desc
        p.font.name = BRAND_BODY_FONT
        p.font.size = Pt(15)
        p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # ===== SLIDE 15: Section Break - Client Communication Hub =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

    num_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(4.0), Inches(2.0))
    p = num_box.text_frame.paragraphs[0]
    p.text = "05"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(3.8), Inches(3.0), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT)
    line.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(10.0), Inches(1.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Client Communication Hub"
    p.font.name = BRAND_HEADING_FONT
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.5), Inches(10.0), Inches(0.8))
    p = sub_box.text_frame.paragraphs[0]
    p.text = "The #1 bar complaint is lack of communication — not bad lawyering"
    p.font.name = BRAND_BODY_FONT
    p.font.size = Pt(20)
    p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    return prs

if __name__ == "__main__":
    prs = create_batch3()
    prs.save("part3.pptx")
    print("✓ Batch 3 complete: Slides 11-15 saved to part3.pptx")
